from __future__ import annotations

from pathlib import Path
import shutil

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from matplotlib import font_manager as fm
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
JH_DIR = ROOT / "presentation" / "jh"
ASSET_DIR = JH_DIR / "assets"
FIGURE_PATH = ASSET_DIR / "rq_conclusion_charts.png"

EFFORT_SUMMARY = ROOT / "eval" / "results" / "combined_effort_summary.md"
ABLATION_SUMMARY = ROOT / "eval" / "results" / "combined_ablation_summary.md"

NAVY = "#163866"
BLUE = "#2F6FAF"
GRAY = "#9CA3AF"
LGREEN = "#8FA38A"
BROWN = "#8B5A3C"
DGRAY = "#2D3338"
LGRID = "#E2E4E8"
CREAM = "#F5F3EF"


def setup_font() -> None:
    for candidate in ("Apple SD Gothic Neo", "AppleGothic", "Arial Unicode MS"):
        if any(f.name == candidate for f in fm.fontManager.ttflist):
            plt.rcParams["font.family"] = candidate
            break
    plt.rcParams["axes.unicode_minus"] = False


def find_v1_deck() -> Path:
    matches = sorted(
        p for p in JH_DIR.glob("*v1.pptx") if not p.name.startswith("~$")
    )
    if len(matches) != 1:
        raise RuntimeError(f"Expected one v1 deck, found {len(matches)}: {matches}")
    return matches[0]


def split_md_row(line: str) -> list[str]:
    return [part.strip() for part in line.strip().strip("|").split("|") if part.strip()]


def parse_metric_row(path: Path, row_name: str) -> list[float]:
    for line in path.read_text().splitlines():
        if line.startswith(f"| {row_name} "):
            values = split_md_row(line)[1:]
            return [float(value.replace("+", "").replace("%", "")) for value in values]
    raise RuntimeError(f"Could not find row {row_name!r} in {path}")


def style_axis(ax: plt.Axes, grid_axis: str = "y") -> None:
    ax.set_facecolor("white")
    for side in ("top", "right"):
        ax.spines[side].set_visible(False)
    for side in ("left", "bottom"):
        ax.spines[side].set_color("#C7CBD1")
    ax.tick_params(colors=DGRAY, length=0, labelsize=9)
    if grid_axis == "x":
        ax.xaxis.grid(True, color=LGRID, lw=1)
    else:
        ax.yaxis.grid(True, color=LGRID, lw=1)
    ax.set_axisbelow(True)


def create_chart_figure(path: Path) -> None:
    setup_font()
    path.parent.mkdir(parents=True, exist_ok=True)

    effort_labels = ["low", "medium", "high", "xhigh"]
    effort_diff = parse_metric_row(EFFORT_SUMMARY, "평균 점수 차 (B-A)")

    ablation_diff = parse_metric_row(ABLATION_SUMMARY, "평균 점수 차 (B-A)")
    ablation_low = parse_metric_row(ABLATION_SUMMARY, "95% CI 하한")
    ablation_high = parse_metric_row(ABLATION_SUMMARY, "95% CI 상한")
    overall = ablation_diff[-1]
    overall_ci = (ablation_low[-1], ablation_high[-1])

    stage_labels = ["1", "1~2", "1~3", "전체"]
    stage_ours = np.array([4.84, 4.83, 4.83, 4.81])
    stage_simple = stage_ours - np.array(ablation_diff)

    fig = plt.figure(figsize=(13.35, 3.03), dpi=200, facecolor="white")

    ax1 = fig.add_axes([0.055, 0.18, 0.245, 0.68])
    ax2 = fig.add_axes([0.375, 0.18, 0.285, 0.68])
    ax3 = fig.add_axes([0.735, 0.18, 0.225, 0.68])

    # 1) Overall improvement over Simple.
    ax1.barh([0], [overall], height=0.38, color=NAVY, zorder=3)
    ax1.errorbar(
        overall,
        0,
        xerr=[[overall - overall_ci[0]], [overall_ci[1] - overall]],
        fmt="none",
        ecolor=DGRAY,
        elinewidth=1.5,
        capsize=4,
        zorder=4,
    )
    ax1.axvline(0, color="#C7CBD1", lw=1)
    ax1.text(overall_ci[1] + 0.035, 0.14, f"+{overall:.3f}", va="center", ha="left", fontsize=13, fontweight="bold", color=NAVY)
    ax1.text(0.02, -0.42, "23승 / 1패", va="center", ha="left", fontsize=9.5, color=DGRAY)
    ax1.set_xlim(0, 1.14)
    ax1.set_ylim(-0.62, 0.55)
    ax1.set_yticks([0])
    ax1.set_yticklabels(["전체"])
    ax1.set_xlabel("평균 점수차", fontsize=9.5, color=DGRAY, labelpad=3)
    style_axis(ax1, grid_axis="x")

    # 2) Stage-wise Simple vs Ours scores.
    x = np.arange(len(stage_labels))
    width = 0.34
    ax2.bar(x - width / 2, stage_simple, width=width, color=GRAY, zorder=3, label="Simple")
    ax2.bar(x + width / 2, stage_ours, width=width, color=NAVY, zorder=3, label="Ours")
    ax2.set_xticks(x)
    ax2.set_xticklabels(stage_labels, fontsize=9.2)
    ax2.set_ylim(3.95, 5.05)
    ax2.set_yticks([4.0, 4.3, 4.6, 4.9, 5.0])
    ax2.set_ylabel("평균 점수", fontsize=9.5, color=DGRAY, labelpad=4)
    ax2.set_xlabel("Stage 조건", fontsize=9.5, color=DGRAY, labelpad=3)
    ax2.legend(loc="upper center", bbox_to_anchor=(0.56, 1.08), ncol=2, frameon=False, fontsize=8.6, handlelength=1.4, columnspacing=1.2)
    style_axis(ax2)

    # 3) Improvement by reasoning effort.
    x = np.arange(len(effort_labels))
    bars = ax3.bar(x, effort_diff, width=0.56, color=NAVY, zorder=3)
    for bar, value in zip(bars, effort_diff):
        ax3.text(
            bar.get_x() + bar.get_width() / 2,
            value + 0.022,
            f"+{value:.3f}",
            ha="center",
            va="bottom",
            fontsize=9.4,
            fontweight="bold",
            color=NAVY,
        )
    ax3.set_xticks(x)
    ax3.set_xticklabels([label.upper() for label in effort_labels], fontsize=8.8)
    ax3.set_ylim(0, 0.70)
    ax3.set_yticks([0, 0.2, 0.4, 0.6])
    ax3.set_ylabel("평균 점수차", fontsize=9.5, color=DGRAY, labelpad=4)
    ax3.set_xlabel("Reasoning effort", fontsize=9.5, color=DGRAY, labelpad=3)
    style_axis(ax3)

    fig.savefig(path, dpi=200, facecolor="white")
    plt.close(fig)


def replace_conclusion_box(deck_path: Path, output_path: Path, figure_path: Path) -> None:
    shutil.copy2(deck_path, output_path)
    prs = Presentation(output_path)
    slide = prs.slides[7]

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape.text.startswith(("단계형 파이프라인", "품질 개선이", "모델 성능이")):
            shape.width = prs.slide_width - shape.left - 450000

    target = None
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and "각 RQ 별 결론 정리" in shape.text:
            target = shape
            break
    if target is None:
        raise RuntimeError("Could not find the slide 8 conclusion placeholder.")

    left, top, width, height = target.left, target.top, target.width, target.height
    target._element.getparent().remove(target._element)
    slide.shapes.add_picture(str(figure_path), left, top, width=width, height=height)
    prs.save(output_path)


def main() -> None:
    source = find_v1_deck()
    output = source.with_name(f"{source.stem}_rq_charts.pptx")
    create_chart_figure(FIGURE_PATH)
    replace_conclusion_box(source, output, FIGURE_PATH)
    print(FIGURE_PATH)
    print(output)


if __name__ == "__main__":
    main()
