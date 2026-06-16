from __future__ import annotations

from pathlib import Path
import shutil

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
JH_DIR = ROOT / "presentation" / "jh"
ASSET_DIR = JH_DIR / "assets"
FIGURE_PATH = ASSET_DIR / "rq_conclusion_figure.png"


def find_v1_deck() -> Path:
    matches = sorted(
        p for p in JH_DIR.glob("*v1.pptx") if not p.name.startswith("~$")
    )
    if len(matches) != 1:
        raise RuntimeError(f"Expected one v1 deck, found {len(matches)}: {matches}")
    return matches[0]


def font(size: int) -> ImageFont.FreeTypeFont:
    candidates = [
        Path("/System/Library/Fonts/AppleSDGothicNeo.ttc"),
        Path("/System/Library/Fonts/Supplemental/Arial Unicode.ttf"),
        Path("/Library/Fonts/Arial Unicode.ttf"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size=size)
    return ImageFont.load_default()


def text_size(draw: ImageDraw.ImageDraw, text: str, face: ImageFont.ImageFont) -> tuple[int, int]:
    bbox = draw.textbbox((0, 0), text, font=face)
    return bbox[2] - bbox[0], bbox[3] - bbox[1]


def centered_text(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    text: str,
    face: ImageFont.ImageFont,
    fill: str,
) -> None:
    left, top, right, bottom = box
    width, height = text_size(draw, text, face)
    draw.text(
        (left + (right - left - width) / 2, top + (bottom - top - height) / 2 - 2),
        text,
        font=face,
        fill=fill,
    )


def draw_wrapped(
    draw: ImageDraw.ImageDraw,
    text: str,
    xy: tuple[int, int],
    max_width: int,
    face: ImageFont.ImageFont,
    fill: str,
    line_gap: int = 6,
) -> int:
    words = text.split(" ")
    lines: list[str] = []
    current = ""
    for word in words:
        candidate = word if not current else f"{current} {word}"
        if text_size(draw, candidate, face)[0] <= max_width:
            current = candidate
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)

    x, y = xy
    line_height = text_size(draw, "가", face)[1] + line_gap
    for line in lines:
        draw.text((x, y), line, font=face, fill=fill)
        y += line_height
    return y


def draw_card(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    label: str,
    big: str,
    body: str,
    sub: str,
    accent: str,
) -> None:
    left, top, right, bottom = box
    radius = 24
    draw.rounded_rectangle(box, radius=radius, fill="#FFFFFF", outline="#BED3EA", width=3)

    pill = (left + 34, top + 24, left + 190, top + 66)
    draw.rounded_rectangle(pill, radius=21, fill=accent)
    centered_text(draw, pill, label, font(23), "#FFFFFF")

    big_face = font(68 if len(big) <= 9 else 58)
    draw.text((left + 34, top + 88), big, font=big_face, fill="#173B62")

    y = draw_wrapped(
        draw,
        body,
        (left + 36, top + 178),
        right - left - 72,
        font(32),
        "#1D2B3A",
        line_gap=7,
    )
    draw_wrapped(
        draw,
        sub,
        (left + 36, max(y + 8, top + 258)),
        right - left - 72,
        font(25),
        "#4B5F73",
        line_gap=6,
    )


def create_figure(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)

    width, height = 2670, 606
    image = Image.new("RGB", (width, height), "#1E5A93")
    draw = ImageDraw.Draw(image)

    # Subtle structural bands keep the compact figure legible on a slide.
    draw.rectangle((0, 0, width, height), fill="#1E5A93")
    draw.rectangle((0, 0, width, 96), fill="#184B7B")
    draw.rectangle((0, height - 104, width, height), fill="#EAF2F8")
    draw.line((0, 96, width, 96), fill="#78A9D1", width=3)

    title = "핵심 해석: 단계별 생성의 우위는 단계 수보다 '필수 구조 슬롯'으로 설명된다"
    centered_text(draw, (40, 14, width - 40, 84), title, font(43), "#FFFFFF")

    margin_x = 72
    gap = 36
    card_top = 124
    card_bottom = 478
    card_w = (width - margin_x * 2 - gap * 2) // 3
    cards = [
        (
            margin_x,
            card_top,
            margin_x + card_w,
            card_bottom,
            "관찰",
            "23/24 x 4",
            "각 effort 조건에서 단계별 생성 우위",
            "단계별-일괄 평균 +0.430~+0.587 / 1-5점 척도",
            "#2364A5",
        ),
        (
            margin_x + card_w + gap,
            card_top,
            margin_x + card_w * 2 + gap,
            card_bottom,
            "해석",
            "22~23/24",
            "1단계+합성부터 대부분의 이득 확보",
            "추가 단계는 심화와 잡음을 함께 추가",
            "#2F7D78",
        ),
        (
            margin_x + (card_w + gap) * 2,
            card_top,
            width - margin_x,
            card_bottom,
            "집중 차원",
            "실행 · 측정 · 인계",
            "차이가 가장 선명했던 구조 차원",
            "구체성 · 가능성 · 준비도 중심",
            "#8A5B1B",
        ),
    ]
    for left, top, right, bottom, label, big, body, sub, accent in cards:
        draw_card(draw, (left, top, right, bottom), label, big, body, sub, accent)

    strip_top = height - 88
    draw.text(
        (72, strip_top + 4),
        "남은 공백",
        font=font(29),
        fill="#173B62",
    )
    draw.text(
        (246, strip_top + 4),
        "수치 출처 추적 항목: 일괄 0.08 · 단계별 0.17~0.29",
        font=font(29),
        fill="#1D2B3A",
    )
    draw.text(
        (1530, strip_top + 9),
        "조건: gpt-5.5 · 24쌍/조건(8 scenarios x 3 seeds) · LLM-judge, 단계별 생성 먼저 제시(위치편향 가능)",
        font=font(22),
        fill="#4B5F73",
    )

    image.save(path)


def replace_conclusion_box(deck_path: Path, output_path: Path, figure_path: Path) -> None:
    shutil.copy2(deck_path, output_path)
    prs = Presentation(output_path)
    slide = prs.slides[7]

    target = None
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and "각 RQ 별 결론 정리" in shape.text:
            target = shape
            break
    if target is None:
        raise RuntimeError("Could not find the slide 8 conclusion placeholder.")

    left, top, width, height = target.left, target.top, target.width, target.height
    target._element.getparent().remove(target._element)
    slide.shapes.add_picture(str(figure_path), left, top, width=width, height=height)
    prs.save(output_path)


def main() -> None:
    source = find_v1_deck()
    output = source.with_name(f"{source.stem}_rq_conclusion.pptx")
    create_figure(FIGURE_PATH)
    replace_conclusion_box(source, output, FIGURE_PATH)
    print(FIGURE_PATH)
    print(output)


if __name__ == "__main__":
    main()
