#!/usr/bin/env python3
"""Create a polished copy of the Milemate final deck.

The source deck is kept unchanged. This script copies it, then adds light
native PowerPoint shapes that clarify experimental evidence and slide logic.
"""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "기술혁신이론_기말발표.pptx"
TARGET = ROOT / "기술혁신이론_기말발표_정갈본.pptx"

NAVY = "163866"
INK = "2D3338"
GRAY = "6B7280"
MUTED = "8A92A0"
SAGE = "8FA38A"
BROWN = "8B5A3C"
CREAM = "F5F3EF"
TINT = "EEF2F7"
WHITE = "FFFFFF"
LINE = "D4D8DE"
GREEN_TINT = "EEF5EE"
BLUE_TINT = "EDF3FA"
BROWN_TINT = "F7F0EA"

F_XB = "NanumSquare ExtraBold"
F_B = "나눔스퀘어 Bold"
F_R = "나눔스퀘어"
F_NB = "NanumSquare Bold"


def _set_font(run, font=F_R, size=11, bold=None, color=INK):
    run.font.name = font
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rpr.find(qn(tag))
        if el is None:
            el = OxmlElement(tag)
            rpr.append(el)
        el.set("typeface", font)
    run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def _shape(slide, left, top, width, height, fill=WHITE, line=LINE, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(
        shape_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line:
        shp.line.color.rgb = RGBColor.from_string(line)
        shp.line.width = Pt(0.8)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=11,
    color=INK,
    bold=False,
    font=F_R,
    align="L",
    valign=MSO_ANCHOR.MIDDLE,
    fill=None,
    line=None,
):
    if fill:
        shp = _shape(slide, left, top, width, height, fill=fill, line=line)
    else:
        shp = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    for idx, line_text in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = {"L": PP_ALIGN.LEFT, "C": PP_ALIGN.CENTER, "R": PP_ALIGN.RIGHT}[align]
        p.line_spacing = Pt(size * 1.25)
        run = p.add_run()
        run.text = line_text
        _set_font(run, font=font, size=size, bold=bold, color=color)
    return shp


def add_pill(slide, left, top, width, text, fill=NAVY, color=WHITE, size=9.5):
    return add_text(
        slide,
        left,
        top,
        width,
        0.32,
        text,
        size=size,
        color=color,
        bold=True,
        font=F_B,
        align="C",
        fill=fill,
        line=None,
    )


def add_card(slide, left, top, width, height, title, body, accent=NAVY, fill=WHITE):
    shp = _shape(slide, left, top, width, height, fill=fill, line=LINE)
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.space_after = Pt(3)
    r = p.add_run()
    r.text = title
    _set_font(r, F_XB, 13, True, accent)

    for line in body.split("\n"):
        p = tf.add_paragraph()
        p.line_spacing = Pt(14.8)
        p.space_after = Pt(1)
        r = p.add_run()
        r.text = line
        _set_font(r, F_R, 10.8, None, INK)
    return shp


def add_source(slide, text):
    add_text(slide, 0.58, 7.16, 7.4, 0.18, text, size=7.2, color=MUTED, font=F_R)


def add_rule(slide, left, top, width, color=NAVY, height=0.035):
    _shape(slide, left, top, width, height, fill=color, line=None, radius=False)


def add_arrow(slide, left, top, width, height, text=None, fill=BLUE_TINT, color=NAVY):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(fill)
    shp.line.color.rgb = RGBColor.from_string(color)
    shp.line.width = Pt(0.8)
    shp.shadow.inherit = False
    if text:
        tf = shp.text_frame
        tf.clear()
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.12)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        _set_font(r, F_B, 10, True, color)
    return shp


def polish_research_questions(slide):
    add_rule(slide, 0.58, 5.15, 7.9, SAGE)
    cards = [
        ("핵심 결과 1", "22~23 / 24 승\np<0.0001", NAVY),
        ("핵심 결과 2", "실행·측정·인계\n우위 집중", SAGE),
        ("핵심 결과 3", "effort를 올려도\n격차 유지", BROWN),
    ]
    for i, (title, body, color) in enumerate(cards):
        x = 0.80 + i * 3.08
        add_card(slide, x, 5.35, 2.72, 1.00, title, body, accent=color, fill=WHITE)


def polish_experiment_design(slide):
    add_pill(slide, 0.95, 2.70, 2.15, "동일 입력·모델·예산", fill=NAVY, size=8.9)
    add_arrow(slide, 3.18, 2.75, 0.42, 0.22, fill=BLUE_TINT, color=GRAY)
    add_pill(slide, 3.78, 2.70, 1.75, "Simple 문서", fill=BROWN, size=8.9)
    add_pill(slide, 5.88, 2.70, 1.75, "Ours 문서", fill=SAGE, size=8.9)
    add_arrow(slide, 7.82, 2.75, 0.42, 0.22, fill=BLUE_TINT, color=GRAY)
    add_pill(slide, 8.42, 2.70, 1.70, "익명 문서쌍", fill=GRAY, size=8.9)
    add_arrow(slide, 10.32, 2.75, 0.42, 0.22, fill=BLUE_TINT, color=GRAY)
    add_pill(slide, 10.92, 2.70, 1.62, "32항목 판정", fill=NAVY, size=8.9)


def polish_result_slides(prs):
    # Slide 13: effort sweep
    s = prs.slides[12]
    add_pill(s, 9.70, 2.08, 2.35, "effort 효과 아님", fill=BROWN)
    add_source(s, "근거: combined_effort_summary (gpt-5.5, low/medium/high/xhigh)")

    # Slide 14: stage ablation
    s = prs.slides[13]
    add_pill(s, 9.42, 2.02, 2.75, "반직관: 1단계만으로 22승", fill=NAVY)
    add_source(s, "근거: combined_ablation_summary (stages_1~4, medium effort)")

    # Slide 15: dimension wins
    s = prs.slides[14]
    add_pill(s, 9.44, 2.15, 2.65, "실행·측정·인계 우위", fill=SAGE)
    add_source(s, "근거: combined_ablation_summary, 차원별 판정 비중")

    # Slide 16: item dumbbell
    s = prs.slides[15]
    add_text(
        s,
        8.86,
        2.20,
        2.75,
        0.50,
        "읽는 법: 오른쪽으로 갈수록\n항목이 완성됨",
        size=8.8,
        color=GRAY,
        font=F_R,
        align="C",
        fill=WHITE,
        line=LINE,
    )
    add_source(s, "근거: combined_ablation_summary, 항목별 평균 점수")


def polish_stage_contribution(slide):
    examples = [
        (0.88, "실험 예시", "merchant_prep:\nAPI field + KPI owner 완비", SAGE),
        (5.18, "실험 예시", "return_pickup:\ncapacity ledger + hold/confirm", NAVY),
        (9.50, "실험 예시", "cs_triage:\n'미계측/드라이런' 레이블", BROWN),
    ]
    for x, title, body, color in examples:
        add_text(slide, x, 5.25, 3.45, 0.46, f"{title} — {body}", size=8.0, color=color, bold=True, font=F_B, fill=WHITE, line=LINE)
    add_source(slide, "근거: synthesis_report, scenario_data stage transition analysis")


def polish_failure_case(slide):
    chips = [
        ("checkout_fee", "seed1 / stages 1~3"),
        ("rider_onboarding", "seed3 / stages 1~2"),
        ("dispatch", "seed1 / stage 3"),
        ("return_pickup", "seed3 / stage 4"),
    ]
    for i, (name, detail) in enumerate(chips):
        x = 0.78 + i * 3.05
        add_text(
            slide,
            x,
            3.10,
            2.63,
            0.66,
            f"{name}\n{detail}",
            size=8.9,
            color=NAVY if i != 3 else BROWN,
            bold=True,
            font=F_B,
            align="C",
            fill=BLUE_TINT if i != 3 else BROWN_TINT,
            line=None,
        )
    add_rule(slide, 0.78, 3.93, 10.93, SAGE)
    add_source(slide, "근거: synthesis_report 실패 케이스 분석")


def polish_d82_gap(slide):
    add_text(
        slide,
        0.66,
        2.92,
        5.88,
        0.38,
        "실패 패턴: 정밀 수치 -> 출처/가정 없음 -> fail",
        size=9.4,
        color=BROWN,
        bold=True,
        font=F_B,
        align="C",
        fill=BROWN_TINT,
        line=None,
    )
    add_text(
        slide,
        7.15,
        3.36,
        4.72,
        0.46,
        "해결 단서: 수치 삭제가 아니라 '미계측 / 드라이런 조정'으로 라벨링",
        size=9.2,
        color=NAVY,
        bold=True,
        font=F_B,
        align="C",
        fill=GREEN_TINT,
        line=None,
    )
    add_source(slide, "근거: synthesis_report 수치 출처 분석, cs_repeat_inquiry_triage stages_4")


def polish_interpretation(slide):
    cards = [
        ("1", "사람 승인", "단계별 판단을\n명시적으로 남김", NAVY),
        ("2", "구조 슬롯", "KPI owner / 결정 선택지 /\n질문 owner·기한을 강제", SAGE),
        ("3", "개발 인계", "화면 · 정책 · 데이터 ·\n이벤트 로그로 번역", BROWN),
    ]
    for i, (num, title, body, color) in enumerate(cards):
        x = 0.82 + i * 4.08
        add_text(slide, x, 3.82, 0.52, 0.52, num, size=18, color=WHITE, bold=True, font=F_XB, align="C", fill=color, line=None)
        add_card(slide, x + 0.60, 3.70, 3.05, 1.20, title, body, accent=color, fill=WHITE)
        if i < 2:
            add_arrow(slide, x + 3.72, 4.05, 0.55, 0.24, fill=BLUE_TINT, color=GRAY)
    add_text(
        slide,
        0.82,
        5.38,
        11.70,
        0.56,
        "해석: 모델 능력을 키운 것이 아니라, 기획서가 반드시 가져야 할 구조를 빠뜨리지 않게 만든 프로세스 혁신",
        size=12.2,
        color=NAVY,
        bold=True,
        font=F_B,
        align="C",
        fill=BLUE_TINT,
        line=None,
    )


def polish_limitations(slide):
    add_text(
        slide,
        5.88,
        4.10,
        1.58,
        0.46,
        "검증 결과",
        size=10.2,
        color=WHITE,
        bold=True,
        font=F_B,
        align="C",
        fill=NAVY,
        line=None,
    )
    add_arrow(slide, 5.92, 4.68, 1.48, 0.38, "설계 과제", fill=GREEN_TINT, color=SAGE)
    add_text(
        slide,
        4.72,
        5.48,
        3.90,
        0.54,
        "직접 대응: BA/AB 스왑 판정 + 구조 감사 + 수치 출처 패스",
        size=9.3,
        color=NAVY,
        bold=True,
        font=F_B,
        align="C",
        fill=WHITE,
        line=LINE,
    )


def polish_conclusion(slide):
    cards = [
        ("22~23 / 24", "Ours 승\np<0.0001", NAVY),
        ("실행·측정·인계", "구조 차원 압도", SAGE),
        ("수치 출처", "다음 설계 과제\n출처 검증", BROWN),
    ]
    for i, (title, body, color) in enumerate(cards):
        x = 1.15 + i * 3.75
        add_text(slide, x, 5.05, 3.10, 0.58, title, size=18, color=color, bold=True, font=F_XB, align="C", fill=WHITE, line=LINE)
        add_text(slide, x, 5.66, 3.10, 0.64, body, size=10.2, color=INK, bold=True, font=F_B, align="C", fill=CREAM, line=None)
    add_source(slide, "근거: combined_effort_summary, combined_ablation_summary, synthesis_report")


def main():
    if not SOURCE.exists():
        raise FileNotFoundError(SOURCE)
    shutil.copy2(SOURCE, TARGET)

    prs = Presentation(TARGET)
    polish_research_questions(prs.slides[7])
    polish_experiment_design(prs.slides[8])
    polish_result_slides(prs)
    polish_stage_contribution(prs.slides[18])
    polish_failure_case(prs.slides[19])
    polish_d82_gap(prs.slides[20])
    polish_interpretation(prs.slides[21])
    polish_limitations(prs.slides[22])
    polish_conclusion(prs.slides[23])

    prs.save(TARGET)
    print(f"saved {TARGET} slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
