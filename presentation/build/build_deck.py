#!/usr/bin/env python3
"""Build the Milemate final deck (25 slides).
- Keeps recap slides (problem / 3-axis / pipeline / thanks) 100% untouched.
- Authors title, agenda, dividers, and all PART2/PART3 slides.
- Inserts matplotlib charts + custom shape grids + highlighted snippet boxes.
Input : working.pptx (25-slide rearranged base)
Output: ../기술혁신이론_기말발표.pptx
"""
import sys
sys.path.insert(0, "/Users/jaeheemacbook/.claude/skills/office-pptx/scripts")
from pathlib import Path
from inventory import extract_text_inventory
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
import struct

NAVY="163866"; INK="2D3338"; GRAY="6B7280"; BROWN="8B5A3C"; SAGE="8FA38A"
WHITE="FFFFFF"; SUBBLUE="B8C5D6"; MUT="808080"; CREAM="F5F3EF"; TINT="EEF2F7"
F_XB="NanumSquare ExtraBold"; F_B="나눔스퀘어 Bold"; F_R="나눔스퀘어"; F_NB="NanumSquare Bold"
ALIGN={"L":PP_ALIGN.LEFT,"C":PP_ALIGN.CENTER,"R":PP_ALIGN.RIGHT}
REPO_ROOT=Path(__file__).resolve().parents[2]

def P(text,font=F_R,size=None,bold=None,color=INK,bullet=False,level=0,align=None,sb=None,sa=None,ls=None):
    return dict(text=text,font=font,size=size,bold=bold,color=color,bullet=bullet,
                level=level,align=align,sb=sb,sa=sa,ls=ls)

def _set_font(run,font,size,bold,color):
    if font:
        run.font.name=font
        rPr=run._r.get_or_add_rPr()
        for tag in ("a:ea","a:cs"):
            el=rPr.find(qn(tag))
            if el is None:
                el=OxmlElement(tag); rPr.append(el)
            el.set("typeface",font)
    if size is not None: run.font.size=Pt(size)
    if bold is not None: run.font.bold=bold
    if color: run.font.color.rgb=RGBColor.from_string(color)

def _bullet(p,on,level,size):
    pPr=p._p.get_or_add_pPr()
    for c in list(pPr):
        if c.tag.endswith(("buChar","buNone","buAutoNum","buFont")): pPr.remove(c)
    if on:
        p.level=level
        fs=size or 14.0
        pPr.set("marL",str(int((fs*(1.6+level*1.6))*12700)))
        pPr.set("indent",str(int(-fs*0.8*12700)))
        bc=OxmlElement("a:buChar"); bc.set("char","•"); pPr.append(bc)
    else:
        pPr.set("marL","0"); pPr.set("indent","0")
        pPr.insert(0,OxmlElement("a:buNone"))

def fill(shape,paras):
    tf=shape.text_frame; tf.clear()
    for i,d in enumerate(paras):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        _bullet(p,d["bullet"],d["level"],d["size"])
        if d["align"] in ALIGN: p.alignment=ALIGN[d["align"]]
        if d["sb"] is not None: p.space_before=Pt(d["sb"])
        if d["sa"] is not None: p.space_after=Pt(d["sa"])
        if d["ls"] is not None: p.line_spacing=Pt(d["ls"])
        r=p.add_run(); r.text=d["text"]
        _set_font(r,d["font"],d["size"],d["bold"],d["color"])

def sg(text,font=F_R,size=11.5,bold=None,color=INK):
    return (text,font,size,bold,color)
def LN(*segs,**kw):
    return {"segs":list(segs),**kw}
def fill_seg(shape,lines):
    tf=shape.text_frame; tf.clear()
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        _bullet(p,ln.get("bullet",False),ln.get("level",0),ln.get("size",11.5))
        p.alignment=ALIGN.get(ln.get("align","L"),PP_ALIGN.LEFT)
        if ln.get("sb") is not None: p.space_before=Pt(ln["sb"])
        if ln.get("sa") is not None: p.space_after=Pt(ln["sa"])
        if ln.get("ls") is not None: p.line_spacing=Pt(ln["ls"])
        for (txt,font,size,bold,color) in ln["segs"]:
            r=p.add_run(); r.text=txt; _set_font(r,font,size,bold,color)

def edit_run(shape,pidx,text):
    p=shape.text_frame.paragraphs[pidx]
    if p.runs:
        p.runs[0].text=text
        for r in p.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        p.add_run().text=text

def strip_mid_pics(slide,min_top_in=0.9):
    for shp in list(slide.shapes):
        if shp.shape_type==MSO_SHAPE_TYPE.PICTURE and shp.top is not None and shp.top>Inches(min_top_in):
            shp._element.getparent().remove(shp._element)

def png_size(path):
    with open(path,"rb") as f:
        f.read(16); w,h=struct.unpack(">II",f.read(8))
    return w,h
def add_chart(slide,path,width_in,top_in,max_bottom=7.32):
    w,h=png_size(path); ar=h/w
    width=width_in; height=width*ar
    if top_in+height>max_bottom:
        height=max_bottom-top_in; width=height/ar
    left=(13.333-width)/2
    slide.shapes.add_picture(path,Inches(left),Inches(top_in),width=Inches(width))
def add_pic(slide,path,left,top,width):
    slide.shapes.add_picture(path,Inches(left),Inches(top),width=Inches(width))

def add_box(slide,left,top,w,h,fill_hex,line_hex,anchor="t"):
    sp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(left),Inches(top),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=RGBColor.from_string(fill_hex)
    if line_hex: sp.line.color.rgb=RGBColor.from_string(line_hex); sp.line.width=Pt(1)
    else: sp.line.fill.background()
    sp.shadow.inherit=False
    tf=sp.text_frame; tf.word_wrap=True
    tf.vertical_anchor=MSO_ANCHOR.TOP if anchor=="t" else MSO_ANCHOR.MIDDLE
    tf.margin_left=Inches(0.18); tf.margin_right=Inches(0.1); tf.margin_top=Inches(0.14); tf.margin_bottom=Inches(0.06)
    return sp
def add_text(slide,left,top,w,h,paras):
    tb=slide.shapes.add_textbox(Inches(left),Inches(top),Inches(w),Inches(h))
    tb.text_frame.word_wrap=True
    fill(tb,paras); return tb
def pill(slide,left,top,w,h,fill_hex,text,color,size=12,line_hex=None):
    b=add_box(slide,left,top,w,h,fill_hex,line_hex,anchor="m")
    fill(b,[P(text,font=F_B,size=size,bold=True,color=color,align="C")])
    return b
def arrow_down(slide,left,top,w):
    add_text(slide,left,top,w,0.26,[P("↓",font=F_B,size=14,bold=True,color=SAGE,align="C")])

def read_repo_text(rel):
    return (REPO_ROOT/rel).read_text(encoding="utf-8")

def compact_excerpt(text,max_lines=24,max_chars=1700):
    lines=[ln.rstrip() for ln in text.splitlines() if ln.strip()]
    out=[]; chars=0
    for ln in lines:
        if len(out)>=max_lines or chars+len(ln)>max_chars:
            break
        out.append(ln); chars+=len(ln)+1
    return "\n".join(out)

def set_notes(slide,text):
    ns=slide.notes_slide
    ns.notes_text_frame.text=text

def card(num_title,cap,cap_col,pairs):
    out=[P(num_title,font=F_XB,size=16,bold=True,color=NAVY,sa=1),
         P(cap,font=F_B,size=10,bold=True,color=cap_col,sa=5)]
    for b,g in pairs:
        out.append(P(b,font=F_B,size=12,bold=True,color=INK,sa=0.5,ls=15.6))
        if g:
            out.append(P(g,font=F_R,size=10.5,color=GRAY,sa=3.5,ls=13.6))
    return out

ASSET="../assets"
prs=Presentation("working.pptx")
S=prs.slides
INV=extract_text_inventory(Path("working.pptx"),prs)
def sh(i):
    return [sd.shape for sd in INV[f"slide-{i}"].values()]
def hdr(i,obj,title):
    s=sh(i)
    fill(s[1],[P(obj,font=F_XB,size=14,bold=True,color=INK)])
    fill(s[2],[P(title,font=F_XB,size=28,bold=True,color=NAVY)])
    return s

# ===== slide 0: title =====
edit_run(sh(0)[0],1,"기획자를 위한 기술 기반 기획서 작성 보조 에이전트 — 설계와 검증")

# ===== slide 1: agenda =====
fill(sh(1)[1],[
 P("PART 1.",font=F_XB,size=20,bold=True,color=NAVY,sa=4,ls=30),
 P("시스템 개요 — 중간발표 요약",font=F_R,size=16,color=MUT,sa=4,ls=24),
 P("기획–개발 번역 병목과 문제 인식",font=F_B,size=15,bold=True,bullet=True,level=1,sa=2,ls=22.5),
 P("Milemate의 4단계 / Human-in-the-loop 파이프라인",font=F_B,size=15,bold=True,bullet=True,level=1,sa=2,ls=22.5),
 P("PART 2.",font=F_XB,size=20,bold=True,color=NAVY,sb=6,sa=4,ls=30),
 P("방법론 유효성 검증 — LLM-as-a-judge 비교 (Ours vs Simple)",font=F_R,size=16,color=MUT,sa=4,ls=24),
 P("시나리오/평가 설계 (8차원 32항목 / 공정성/편향 통제)",font=F_B,size=15,bold=True,bullet=True,level=1,sa=2,ls=22.5),
 P("결과/사례: Effort/Ablation / 차원/항목 / 실제 문서 예시",font=F_B,size=15,bold=True,bullet=True,level=1,sa=2,ls=22.5),
 P("PART 3.",font=F_XB,size=20,bold=True,color=NAVY,sb=6,sa=4,ls=30),
 P("통찰과 결론",font=F_R,size=16,color=MUT,sa=4,ls=24),
 P("단계별 기여 / 실패 케이스 / 구조적 공백",font=F_B,size=15,bold=True,bullet=True,level=1,sa=2,ls=22.5),
 P("종합 해석 / 한계 / 결론",font=F_B,size=15,bold=True,bullet=True,level=1,sa=2,ls=22.5),
])

# ===== slide 2: PART1 divider =====
s=sh(2); edit_run(s[2],0,"시스템 개요")
edit_run(s[3],0,"중간발표 요약 — 문제 인식 / 제안 방법론 / 에이전트 파이프라인")

# ===== slide 3: 문제 (KEEP) — typo fix only =====
for shp in sh(3):
    if "기술 제약 지식" in shp.text_frame.text:
        for p in shp.text_frame.paragraphs:
            if p.runs and "기술 제약 지식" in "".join(r.text for r in p.runs):
                p.runs[0].text="아이디어는 있지만 기술 제약 지식 부족"
                for r in p.runs[1:]:
                    r._r.getparent().remove(r._r)
        break

# ===== slide 4 (3-axis) & 5 (pipeline): KEEP =====

# ===== slide 6: PART2 divider =====
s=sh(6); edit_run(s[2],0,"방법론 유효성 검증")
edit_run(s[3],0,"단계형 파이프라인은 정말 더 나은 기획서를 만드는가? — 동일 입력 Ours vs Simple 비교")

# ===== slide 7: 연구 질문 =====
strip_mid_pics(S[7])
s=hdr(7,"04. 무엇을 검증하는가","연구 질문")
fill(s[3],[
 P("중간발표는 효용을 \"주장\"했다. 이번에는 동일 입력으로 단계형(Ours)과 일회성(Simple)을 비교해 검증한다.",font=F_R,size=15,color=INK,sa=10,ls=22),
 P("RQ1.  단계형 파이프라인(Ours)이 일회성 생성(Simple)보다 더 나은 기획서를 만드는가?",font=F_NB,size=15,bold=True,color=NAVY,sa=8,ls=23),
 P("RQ2.  그 이득은 어느 단계 / 어떤 속성에서 오는가?  (단계를 늘릴수록 좋아지는가?)",font=F_NB,size=15,bold=True,color=NAVY,sa=8,ls=23),
 P("RQ3.  결과가 모델의 reasoning effort에 좌우되는가?  (effort만 높이면 되는 것 아닌가?)",font=F_NB,size=15,bold=True,color=NAVY,sa=8,ls=23),
 P("→ RQ2가 이 발표의 핵심 답을 준다.",font=F_B,size=15,bold=True,color=INK,sb=6,ls=22),
])

# ===== slide 8: 실험 설계 (Simple vs Ours, 2-card) =====
s=hdr(8,"05. 평가 설계 / Simple vs Ours","실험 설계 — 같은 입력, 두 경로")
fill(s[3],[P("동일 initial_input / 동일 생성 모델 / 동일 토큰 예산.  8개 라스트마일 시나리오 × seed ≥ 3쌍.",font=F_R,size=14,color=INK,ls=19)])
fill(s[4],[
 P("Simple / 일회성 생성",font=F_XB,size=18,bold=True,color=NAVY,sa=4),
 P("ONE-SHOT (baseline)",font=F_R,size=10,color=BROWN,sa=9),
 P("중간 단계 없이 한 번에 최종 기획서 생성",font=F_R,size=13,bullet=True,sa=4,ls=18),
 P("일반적인 LLM 사용 방식",font=F_R,size=13,bullet=True,ls=18),
])
fill(s[5],[
 P("Ours / 단계형 파이프라인",font=F_XB,size=18,bold=True,color=NAVY,sa=4),
 P("STAGED PIPELINE (Milemate)",font=F_R,size=10,color=SAGE,sa=9),
 P("문제정의 → MVP범위 → 리스크검증 → 최종",font=F_R,size=13,bullet=True,sa=4,ls=18),
 P("각 단계 승인 / rollback 후 종합 합성",font=F_R,size=13,bullet=True,ls=18),
])
fill(s[6],[P("공정성: 모델/입력/토큰을 동일하게 고정 — 차이는 오직 \"단계 구조\"의 유무",font=F_B,size=16,bold=True,color=INK,ls=22)])

# ===== slide 9: 시나리오 소개 (8개 목록 + 사용자 첫 질문) =====
strip_mid_pics(S[9]); s=hdr(9,"06. 평가 시나리오","현업에서 마주칠 8개 라스트마일 시나리오")
fill(s[3],[P("단일 사례가 아니라 성격이 다른 8개 운영 문제를 직접 설계 — 신규 도입/기존 개선/규제 민감을 의도적으로 포괄.",font=F_NB,size=13.5,bold=True,color=NAVY,ls=19)])
add_text(S[9],0.55,2.4,6.45,4.6,[
 P("8개 시나리오",font=F_XB,size=13,bold=True,color=NAVY,sa=6),
 P("동적 배차/경로 재추천 — 피크 배차 집중 완화",font=F_R,size=12.5,bullet=True,sa=4.5,ls=17),
 P("ETA 예측/지연 알림 — 지연 주문 선제 감지/안내",font=F_R,size=12.5,bullet=True,sa=4.5,ls=17),
 P("실패 배송 리스크/사전 개입 — 실패 전 미리 조치",font=F_R,size=12.5,bullet=True,sa=4.5,ls=17),
 P("신규 라이더 온보딩/이탈 방지 — 첫 2주 선제 연락",font=F_R,size=12.5,bullet=True,sa=4.5,ls=17),
 P("반품/교환 회수 예약/추적 — 역배송 흐름 신설",font=F_R,size=12.5,bullet=True,sa=4.5,ls=17),
 P("결제 전 배송비 투명 안내 — 장바구니 이탈 감소",font=F_R,size=12.5,bullet=True,sa=4.5,ls=17),
 P("매장 준비시간 지연 가시화 — 기사 헛대기 감소",font=F_R,size=12.5,bullet=True,sa=4.5,ls=17),
 P("반복 배송 문의 자동 분류 — 단순 문의 1차 자동",font=F_R,size=12.5,bullet=True,ls=17),
])
qb=add_box(S[9],7.25,2.4,5.55,4.6,CREAM,"DCDCDC")
fill_seg(qb,[
 LN(sg("사용자의 첫 질문",F_XB,14,True,NAVY),sa=2),
 LN(sg("예: 동적 배차 / 강남/서초 저녁 피크",F_R,10.5,None,BROWN),sa=8),
 LN(sg("“저녁 6~8시 강남/서초에 주문이 몰리면 잘하는 기사 몇 분께만 콜이 쏠려, 늦은 주문 하나에 사과 쿠폰 3천 원씩 나갑니다.",F_R,12,None,INK),sa=5,ls=18),
 LN(sg("지금은 운영팀이 주문/기사 위치 화면을 눈으로 번갈아 보며 손으로 바꾸는데 피크엔 손이 안 따라가요.",F_R,12,None,INK),sa=5,ls=18),
 LN(sg("위험한 주문이 먼저 뜨고 왜 위험한지 이유까지 보여 승인만 누르면 좋겠는데 — ",F_R,12,None,INK),sg("개발은 1도 몰라서 뭐부터 말해야 할지조차 모르겠어요.”",F_NB,12,True,NAVY),sa=9,ls=18),
 LN(sg("→ 이 자연어 한 문단이 Stage 1 입력 (initial_input)",F_B,11.5,True,SAGE)),
])

# ===== slide 10: 평가 척도 (2-card) =====
s=hdr(10,"07. 평가 척도 / 8차원 32항목","평가 척도 — LLM-as-a-judge")
fill(s[3],[P("두 문서를 익명(문서1/2)으로 제시, 항목마다 pass(1.0)/partial(0.5)/fail(0)/원문 인용 필수. 차원 점수는 파이프라인이 산식으로 계산.",font=F_R,size=14,color=INK,ls=19)])
fill(s[4],[
 P("판정 방식",font=F_XB,size=18,bold=True,color=NAVY,sa=4),
 P("HOW",font=F_R,size=10,color=BROWN,sa=9),
 P("32개 원자 항목을 기계적으로 채점",font=F_R,size=13,bullet=True,sa=4,ls=18),
 P("pass/partial은 반드시 원문 인용",font=F_R,size=13,bullet=True,sa=4,ls=18),
 P("차원 = 1 + 4 × (획득 / 유효항목)",font=F_R,size=13,bullet=True,ls=18),
])
fill(s[5],[
 P("8개 차원",font=F_XB,size=18,bold=True,color=NAVY,sa=4),
 P("WHAT",font=F_R,size=10,color=SAGE,sa=9),
 P("문제정의 / 범위규율 / 실행구체성 / 측정가능성",font=F_R,size=13,bullet=True,sa=4,ls=18),
 P("리스크·보류 / 내적일관성 / 인계준비 / 입력충실성",font=F_R,size=13,bullet=True,ls=18),
])
fill(s[6],[P("핵심 원칙: 산출 형식(섹션/용어)이 아니라 \"내용 속성\"만 측정 → 순환 논증 회피",font=F_B,size=16,bold=True,color=INK,ls=22)])

# ===== slide 11: 8차원을 도형으로 =====
strip_mid_pics(S[11]); s=hdr(11,"07. 평가 척도 / 8차원","무엇을 보는가 — 8개 평가 차원")
fill(s[3],[P("각 차원은 \"좋은 기획서가 갖춰야 할 내용 속성\"이며, 아래 32개 원자 항목으로 채점된다.",font=F_R,size=13,color=INK,ls=18)])
DIMS=[("문제정의 품질","고통 장면 구체성 / 임팩트 수치화 / 현행 우회/한계 / 성공기준 검증가능성"),
      ("범위 규율","포함 범위 구체성 / 제외 범위 명시 / 자동화–사람 경계 / 범위 결정 근거"),
      ("실행 구체성","화면/접점 정의 / 예외 상태 처리 / 운영정책 실행가능성 / 데이터 요건 특정 / 측정용 로그/이벤트"),
      ("측정 가능성","KPI baseline/target / 측정 방법 명시 / 측정 담당 주체 / 문제–KPI 정합성"),
      ("리스크/보류","시나리오 고유 실패 경로 / 보류/중단 발동 조건 / 과잉 자동화 가드레일 / 개인정보/민감정보 취급"),
      ("내적 일관성","문제→범위→지표 흐름 단절 / 문서 내 모순 / 정책–설계 충돌"),
      ("의사결정/인계","결정 안건 완결성 / 미해결 질문 추적 / 개발팀 인계 질문/요청 / 결재 판단 완결성"),
      ("입력 충실/절제","화자 우려/제약 반영 / 수치 출처 추적 / 요청 방향 유지 / 효과 주장 절제")]
bw,bh,x0,y0,dx,dy=5.9,1.1,0.55,2.4,6.2,1.24
for idx,(name,desc) in enumerate(DIMS):
    col,row=idx%2,idx//2
    box=add_box(S[11],x0+col*dx,y0+row*dy,bw,bh,CREAM,"DCDCDC",anchor="m")
    fill_seg(box,[
     LN(sg(name,F_B,13.2,True,NAVY),sa=3),
     LN(sg(desc,F_R,10.5,None,GRAY),ls=13.5),
    ])

# ===== slide 12: 결과 실험1 (effort grouped bars) =====
strip_mid_pics(S[12]); s=hdr(12,"08. 결과 / 실험 1","Effort를 올려도 Ours 우위는 그대로")
fill(s[3],[P("각 effort에서 Simple/Ours를 묶어 비교해도 격차는 유지된다.",font=F_R,size=13.5,color=INK,ls=18)])
add_chart(S[12],f"{ASSET}/fig_effort.png",10.6,2.35)

# ===== slide 13: 결과 실험2 (ablation) =====
strip_mid_pics(S[13]); s=hdr(13,"09. 결과 / 실험 2","단계별 생성은 일괄 생성을 앞선다")
fill(s[3],[P("각 stage 조건에서 단계별 생성(B)을 일괄 생성(A)과 직접 비교했다. 축 하한을 올려 성능 차이를 강조했다.",font=F_R,size=13.5,color=INK,ls=18)])
add_chart(S[13],f"{ASSET}/fig_ablation.png",9.7,2.25)

# ===== slide 14: 결과 차원별 =====
strip_mid_pics(S[14]); s=hdr(14,"10. 결과 / 차원별","실행·측정·인계에서 압도"); fill(s[3],[])
add_chart(S[14],f"{ASSET}/fig_dimension.png",9.9,2.25)

# ===== slide 15: 결과 항목별 (dumbbell) =====
strip_mid_pics(S[15]); s=hdr(15,"11. 결과 / 항목별","방법론이 정확히 무엇을 채우는가"); fill(s[3],[])
add_chart(S[15],f"{ASSET}/fig_dumbbell.png",9.1,1.95)

# ===== slide 16: 예시① 실제 출력문 비교 =====
strip_mid_pics(S[16]); s=hdr(16,"12. 예시로 보는 차이 ①","Simple/Ours 실제 출력문")
fill(s[3],[P("동일한 사용자 질문에서 나온 두 출력물을 문서 카드로 비교한다. 전체 원문은 이 슬라이드 노트에 보존.",font=F_R,size=13,color=INK,ls=18)])
simple_out=read_repo_text("eval/results/smoke1/generations/dispatch_recommendation/seed1/A.md")
ours_out=read_repo_text("eval/results/smoke1/generations/dispatch_recommendation/seed1/B.md")
qb=add_box(S[16],0.55,2.25,12.25,0.88,CREAM,"DCDCDC",anchor="m")
fill_seg(qb,[
 LN(sg("사용자 질문 ",F_B,11,True,BROWN),sg("(동적 배차)",F_R,10,None,GRAY),sa=2),
 LN(sg("“위험한 주문이 먼저 뜨고, 왜 위험한지 이유까지 보여 운영자가 승인만 누르면 좋겠어요. 그런데 개발은 1도 몰라서 뭐부터 말해야 할지 모르겠어요.”",F_R,11.8,None,INK),ls=16),
])

def output_card(slide,left,title,meta,text,accent):
    card=add_box(slide,left,3.35,5.95,3.35,WHITE,"DCDCDC")
    fill_seg(card,[
     LN(sg(title,F_XB,14.5,True,accent),sa=1),
     LN(sg(meta,F_R,8.8,None,GRAY),sa=5),
     LN(sg(compact_excerpt(text,max_lines=14,max_chars=1050),F_R,6.1,None,INK),ls=7.0),
    ])
    return card
output_card(S[16],0.55,"Simple / 일회성 출력","A.md · 150 lines · 최종 문서형",simple_out,BROWN)
output_card(S[16],6.85,"Ours / 단계형 출력","B.md · 1,247 lines · 구조화 산출물",ours_out,NAVY)
add_text(S[16],0.55,6.86,12.25,0.34,[P("차이의 핵심: Simple은 읽기 쉬운 최종 문서에 가깝고, Ours는 결정·KPI·데이터·보류 조건을 구조화해 인계 가능한 형태로 만든다.",font=F_B,size=11.2,bold=True,color=INK,ls=15)])
set_notes(S[16],
          "[전체 출력물 - Simple / A.md]\n\n"
          + simple_out
          + "\n\n\n[전체 출력물 - Ours / B.md]\n\n"
          + ours_out)

# ===== slide 17: 예시② 바람직한 조건 + 운영 구조 =====
strip_mid_pics(S[17]); s=hdr(17,"13. 예시로 보는 차이 ②","바람직한 조건이 운영 구조로 바뀐다")
fill(s[3],[P("바람직한 출력 조건을 위에 고정하고, 아래에는 Ours가 실제로 만든 운영 구조 예시를 배치했다.",font=F_R,size=13,color=INK,ls=18)])
CONDS=["KPI 담당자(owner)", "결정 선택지 2개+", "질문 책임자·기한", "정책 trigger/rule/owner/exception"]
for i,label in enumerate(CONDS):
    x=0.64+i*3.1
    pill(S[17],x,2.35,2.75,0.48,TINT,label,NAVY,size=10.2,line_hex=NAVY)

def case_col(left,header,stage,q,cond,steps):
    add_text(S[17],left,3.02,3.85,0.45,[P(header,font=F_XB,size=12.8,bold=True,color=NAVY,sa=0),P(stage,font=F_R,size=8.8,color=SAGE)])
    add_text(S[17],left,3.58,3.85,0.72,[P("질문",font=F_B,size=9,bold=True,color=BROWN,sa=0.5),P(q,font=F_R,size=9.6,color=INK,ls=12.8)])
    add_text(S[17],left,4.38,3.85,0.62,[P("충족 조건",font=F_B,size=9,bold=True,color=SAGE,sa=0.5),P(cond,font=F_R,size=9.6,color=INK,ls=12.8)])
    fy=5.10
    for i,st in enumerate(steps):
        pill(S[17],left+0.25,fy,3.35,0.42,TINT,st,NAVY,size=9.6,line_hex=NAVY); fy+=0.42
        if i<len(steps)-1:
            arrow_down(S[17],left+0.25,fy-0.03,3.35); fy+=0.22
case_col(0.55,"반품 회수 예약","Stage 3 / 타당성 검증","“물건을 언제 가지러 오는지 정해진 게 없어요”","초과 예약 없이 안전하게 자리를 잡아야 한다",
 ["예약 요청","임시 보류","확정 / 자동 해제"])
case_col(4.6,"실패 배송 개입","Stage 3 / 타당성 검증","“실패가 터진 다음에야 수습을 시작해요”","여러 신호로 개입 수준을 미리 정해야 한다",
 ["위험 신호 종합","개입 등급 산정","5단계 조치"])
case_col(8.65,"동적 배차 인계","Stage 4 / 최종화","“개발팀에 뭐부터 말해야 할지 모르겠어요”","개발이 바로 연동할 인터페이스가 있어야 한다",
 ["추천 결과","이유 / 신선도 / 비용위험","개발 연동"])

# ===== slide 18: 단계별 기여 (3-card) =====
s=hdr(18,"14. 단계별 기여","각 단계가 실제로 더하는 것")
fill(s[3],[P("단계 추가 = 순이득이 아니다. 후기 단계일수록 정교화와 잡음이 함께 늘어난다.",font=F_R,size=14,color=INK,ls=19)])
fill(s[4],card("2단계 / 서비스 설계","안정적 이득",BROWN,[
 ("정책 4요소/데이터 필드","실행 구체성 구조 완성"),
 ("KPI 구조화","측정 기준/측정 방식/owner"),
 ("시나리오 의존성 낮음","가장 일관적인 단계")]))
fill(s[5],card("3단계 / 타당성 검증","신규성 ↔ 혼재",BROWN,[
 ("개념적 깊이 추가","개입등급/shadow mode"),
 ("가드레일 정교화","수치 기반 중단 조건"),
 ("동시에 잡음 유입","근거없는 수치/비일관")]))
fill(s[6],card("4단계 / 최종화","시나리오 의존",SAGE,[
 ("인계 자료 강화","인식론적 레이블로 일부 해결"),
 ("때때로 퇴행","문서 목적 어긋나면 하락"),
 ("방향이 갈림","가장 이질적인 전환")]))
add_text(S[18],0.55,5.95,12.25,1.15,[
 P("예시 — 정책 1건이 단계마다 어떻게 깊어지는가 (동적 배차)",font=F_NB,size=12.5,bold=True,color=NAVY,sa=4,ls=17),
 P("Simple \"지연 시 고객 안내\"  →  2단계 trigger/rule/owner/exception 4요소  →  3단계 수치 기반 중단 조건  →  4단계 상태기계/API 계약",font=F_R,size=12,color=INK,ls=17.5),
])

# ===== slide 19: 실패 케이스 (+ 구체 사례) =====
strip_mid_pics(S[19]); s=hdr(19,"15. 실패 케이스 분석","Simple이 이긴 4건의 공통 원인")
fill(s[3],[
 P("4개 패배(checkout_fee / rider_onboarding / dispatch / return_pickup)는 단 하나의 원인을 공유한다.",font=F_NB,size=14,bold=True,color=NAVY,sa=5,ls=19),
 P("Ours가 개념 레이어를 쌓으며 기존 구조 슬롯(KPI owner / 결정 선택지 2개+ / 질문 기한)을 비워두고 넘어감.",font=F_R,size=13.5,color=INK,sa=4,ls=19),
 P("체크리스트는 미완성 구조를 partial이 아니라 hard fail로 처리 → 개념적으로 정교해도 진다.",font=F_R,size=13.5,color=INK,ls=19),
])
cb=add_box(S[19],0.55,3.95,12.25,2.5,CREAM,"DCDCDC")
fill_seg(cb,[
 LN(sg("구체 사례 — checkout_fee / seed1  (stages 1~3, 3개 조건 전패)",F_XB,13,True,NAVY),sa=7),
 LN(sg("Simple:  ",F_NB,12.5,True,BROWN),sg("결정 선택지 2개+ / KPI owner / 질문 기한을 모두 충족 → pass",F_R,12.5,None,INK),sa=5,ls=18),
 LN(sg("Ours:  ",F_NB,12.5,True,NAVY),sg("위험주문 한정(1단계) → quote_id 결제 정합(2단계) → fee_exposure_mode 신뢰도 3-tier 공개(3단계)",F_R,12.5,None,INK),sa=3,ls=18),
 LN(sg("          개념은 단계마다 깊어졌지만 ",F_R,12.5,None,GRAY),sg("KPI 담당자 / 결정 선택지 / 질문 기한을 비워둠 → hard fail",F_NB,12.5,True,NAVY),ls=18),
])
add_text(S[19],0.55,6.65,12.25,0.5,[P("함의: 새 개념을 더할 때 기존 슬롯 충족을 강제하는 \"구조 완전성 감사\" 단계가 필요하다.",font=F_B,size=13,bold=True,color=INK,ls=18)])

# ===== slide 20: 수치 출처 구조적 공백 (+ 구체 예시) =====
strip_mid_pics(S[20]); s=hdr(20,"16. 구조적 공백","수치 출처 — 모든 단계가 못 푸는 공백")
fill(s[3],[P("각 단계가 더하는 정밀 수치(임계값/샘플 하한/차단기)가 출처/가정 없이 등장 — Simple/Ours 공통 약점.",font=F_R,size=13.5,color=INK,ls=18)])
add_pic(S[20],f"{ASSET}/fig_d82.png",0.5,3.3,6.3)
xb=add_box(S[20],7.15,3.3,5.6,3.25,CREAM,"DCDCDC")
fill_seg(xb,[
 LN(sg("근거 없는 수치 — 실제 출현 예",F_XB,13,True,NAVY),sa=7),
 LN(sg("dispatch:  \"30초 이내 노출\"",F_R,12,None,INK),sa=4,ls=17),
 LN(sg("eta:  \"신뢰도 0.65\" / \"오경보율 15%\"",F_R,12,None,INK),sa=4,ls=17),
 LN(sg("checkout:  \"predicted_delta_abs ≤ 500원\"",F_R,12,None,INK),sa=4,ls=17),
 LN(sg("merchant_prep:  4개 조건 전부 fail",F_R,12,None,INK),sa=8,ls=17),
 LN(sg("→ 모두 입력/가정/측정 계획 없이 등장",F_B,12,True,BROWN),sa=6),
 LN(sg("유일 해결(cs_triage 4단계): ",F_R,11.5,None,GRAY),sg("\"미계측/드라이런 보정\" 레이블",F_NB,11.5,True,SAGE),ls=16),
])

# ===== slide 21: 종합 해석 =====
strip_mid_pics(S[21]); s=hdr(21,"17. 종합 해석","가치는 \"단계 수\"가 아니라 \"구조 강제\"에 있다")
fill(s[3],[
 P("정량/정성 종합: 일회성 LLM이 자연스럽게는 안 만드는 거버넌스/측정/정책 구조를, 단계형 합성이 반강제로 채운다.",font=F_NB,size=14.5,bold=True,color=NAVY,sa=10,ls=20),
 P("기술혁신 관점",font=F_NB,size=14,bold=True,color=INK,sa=5,ls=20),
 P("Human-in-the-loop + 구조 강제 = 도구가 아니라 \"프로세스 혁신\"",font=F_R,size=14,bullet=True,sa=3,ls=19),
 P("기획–개발 번역 병목을 개인 역량 의존 → 재현 가능한 구조로 이동",font=F_R,size=14,bullet=True,sa=3,ls=19),
 P("혁신의 효과는 \"더 많은 자동화\"가 아니라 \"무엇을 반드시 갖추게 강제하느냐\"에서 나온다",font=F_R,size=14,bullet=True,ls=19),
])

# ===== slide 22: 한계 + 로드맵 (2-card) =====
s=hdr(22,"18. 한계와 개선 방향","한계와 다음 단계")
fill(s[3],[P("검증이 드러낸 약점은 그대로 다음 설계 과제가 된다.",font=F_R,size=14,color=INK,ls=19)])
fill(s[4],[
 P("한계",font=F_XB,size=18,bold=True,color=NAVY,sa=4),
 P("LIMITATIONS",font=F_R,size=10,color=BROWN,sa=9),
 P("위치 편향: 제시 순서 고정(Ours 먼저, 유리할 수 있음)",font=F_R,size=12.5,bullet=True,sa=4,ls=17),
 P("단일 판정 모델 / 자기선호 가능",font=F_R,size=12.5,bullet=True,sa=4,ls=17),
 P("LLM-judge ↔ 사람 평가 상관 미검증",font=F_R,size=12.5,bullet=True,sa=4,ls=17),
 P("단일 도메인(라스트마일) 8 시나리오",font=F_R,size=12.5,bullet=True,ls=17),
])
fill(s[5],[
 P("개선 로드맵",font=F_XB,size=18,bold=True,color=NAVY,sa=4),
 P("ROADMAP",font=F_R,size=10,color=SAGE,sa=9),
 P("구조 완전성 감사 단계 추가",font=F_R,size=12.5,bullet=True,sa=4,ls=17),
 P("수치 출처 검증 전용 패스",font=F_R,size=12.5,bullet=True,sa=4,ls=17),
 P("가드레일 중복 제거 / 신규 KPI owner 강제",font=F_R,size=12.5,bullet=True,sa=4,ls=17),
 P("순서 스왑(BA/AB) 2회 판정으로 강화",font=F_R,size=12.5,bullet=True,ls=17),
])
fill(s[6],[P("한계는 곧 로드맵 — 평가가 개선 지점을 정확히 지목한다",font=F_B,size=16,bold=True,color=INK,ls=22)])

# ===== slide 23: 결론 =====
strip_mid_pics(S[23]); s=hdr(23,"정리","결론")
fill(s[3],[
 P("단계형 Human-in-the-loop 기획 에이전트(Ours)는 일회성 생성(Simple)을 유의하게 능가한다 (22~23 / 24, p<0.0001).",font=F_NB,size=16,bold=True,color=NAVY,sa=10,ls=22),
 P("우위의 본질: 단계 수가 아니라 \"구조적 완전성의 강제\"",font=F_R,size=14.5,bullet=True,sa=4,ls=20),
 P("기여 ① 번역 병목을 푸는 시스템   ② 순환 논증을 피한 공정한 평가 방법론",font=F_R,size=14.5,bullet=True,sa=10,ls=20),
 P("Takeaway — AI 기획 보조의 핵심은 더 똑똑한 모델이 아니라, 무엇을 반드시 갖추게 강제하는 구조다.",font=F_B,size=15.5,bold=True,color=NAVY,ls=22),
])

# ===== slide 24: 감사합니다 (KEEP) =====

out="../기술혁신이론_기말발표.pptx"
prs.save(out)
print("saved",out,"slides:",len(prs.slides))
