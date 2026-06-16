"""Replace slide 18 of v1 with the 4-stage x 4-element summary, as native shapes.

Keeps the 4 placeholders (title / page-no / section header / subtitle), rewrites
the header+subtitle text, removes the old example-flow shapes, and draws a
4x4 card grid matching the deck's tokens.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

SRC = "presentation/jh/기술혁신이론_기말발표_v1.pptx"

# tokens
NAVY = RGBColor(0x16, 0x38, 0x66)
INK = RGBColor(0x2D, 0x33, 0x38)
BROWN = RGBColor(0x8B, 0x5A, 0x3C)
SAGE = RGBColor(0x8F, 0xA3, 0x8A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CELLBG = RGBColor(0xFB, 0xFC, 0xFD)
CELLLN = RGBColor(0xCD, 0xD6, 0xE2)
STAGE_BG = [RGBColor(0x16, 0x38, 0x66), RGBColor(0x27, 0x56, 0x8C),
            RGBColor(0x3A, 0x6E, 0xA5), RGBColor(0x4F, 0x86, 0xC6)]
F_XB = "NanumSquare ExtraBold"
F_B = "나눔스퀘어 Bold"
F_R = "나눔스퀘어"

KEEP = {"제목 1", "슬라이드 번호 개체 틀 2", "내용 개체 틀 3", "텍스트 개체 틀 4"}
NEW_HDR = "13. 단계별 핵심 요약"
NEW_SUB = "한 번의 입력이 4단계로 구조화되는 과정"

COLS = [
    ("Reasoning", "그 단계의 판단", NAVY),
    ("기획자 관점", "planner_view", SAGE),
    ("개발자 관점", "engineer_view", BROWN),
    ("PRD 핵심", "prd_packet", GRAY),
]
STAGES = [
    ("1단계", "문제 정의", [
        "‘자동 배차’가 아니라 **위험 주문 우선순위화 + 운영자 승인**으로 문제 재정의",
        "피크 때 **위험 주문만 노출**, 사유·기사 부담 기준 확인 후 승인",
        "위험점수/추천/승인 API 분리 — **추천과 실행 사이 ‘승인’ 필수**",
        "위험 주문만 상단에 올리는 **승인형 배차 추천 MVP** 정의",
    ]),
    ("2단계", "서비스 구조", [
        "완전 자동 배차 제외, **승인·보류·거절** 흐름 + 4대 지표 KPI 통합",
        "확인→사유→**승인/보류/거절**→KPI 리뷰까지 운영 화면 정의",
        "추천 **상태 모델** 정의 — **approved만** 실제 배차 요청 생성",
        "1단계 범위 유지 + 승인·보류·거절 기록과 **지표 리뷰 흐름 확정**",
    ]),
    ("3단계", "검증 / 리스크", [
        "데이터 늦거나 공급 부족이면 **추천 말고 보류·수동확인으로 전환**",
        "**‘추천을 믿어도 되는 조건’** 정의 — 오래된 데이터 시 보류·롤백",
        "승인 직전 **신선도 재검증 게이트**, blocked_for_stale_data 기록",
        "지연·공급부족·수요초과에서 **믿을 추천 vs 보류 기준** 정의",
    ]),
    ("4단계", "최종 기획서", [
        "조기 발견 + 납득 추천으로 **신뢰·부담·클레임·비용 동시 절감**",
        "문제 재정의 — **늦을 주문을 먼저 찾아 근거로 승인**하게",
        "필수 데이터 묶음 + **구현 순서·검증 계획** 확정",
        "신뢰도·보류 기준 위에 **설명가능 추천을 승인**하는 최종 PRD 확정",
    ]),
]


def set_run(r, text, font, size, bold, color):
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    # set east-asian + complex-script typeface so Korean renders in NanumSquare
    rPr = r._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}" + tag.split(":")[1])
        if el is None:
            from pptx.oxml.ns import qn
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)


def add_seg_para(tf, segs, size, base_font, base_color, align=PP_ALIGN.LEFT, first=True):
    """segs: list of (text, bold). Bold -> navy ExtraBold."""
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    for text, bold in segs:
        r = p.add_run()
        if bold:
            set_run(r, text, F_B, size, True, NAVY)
        else:
            set_run(r, text, base_font, size, False, base_color)
    return p


def parse_bold(text):
    segs, i = [], 0
    parts = text.split("**")
    for k, part in enumerate(parts):
        if part:
            segs.append((part, k % 2 == 1))
    return segs


def rrect(slide, l, t, w, h, fill, line, radius=0.06):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    sp.adjustments[0] = radius
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    for m in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{m}", Inches(0.07))
    return sp


def main():
    prs = Presentation(SRC)
    slide = prs.slides[17]

    # 1) rewrite header + subtitle (keep their fonts; set run[0], drop extras)
    for sh in slide.shapes:
        if sh.name == "내용 개체 틀 3" and sh.has_text_frame:
            p = sh.text_frame.paragraphs[0]
            p.runs[0].text = NEW_HDR
            for r in p.runs[1:]:
                r._r.getparent().remove(r._r)
        if sh.name == "텍스트 개체 틀 4" and sh.has_text_frame:
            p = sh.text_frame.paragraphs[0]
            p.runs[0].text = NEW_SUB
            for r in p.runs[1:]:
                r._r.getparent().remove(r._r)

    # 2) remove old content shapes (keep placeholders)
    for sh in list(slide.shapes):
        if sh.name not in KEEP:
            sh._element.getparent().remove(sh._element)

    # 3) geometry
    GL, GT = 0.45, 1.80          # grid left / top
    LBL_W = 1.55                 # stage label column width
    GAPX, GAPY = 0.12, 0.09
    col_w = (12.88 - (GL + LBL_W) - GAPX * 4) / 4
    col_x = [GL + LBL_W + GAPX + i * (col_w + GAPX) for i in range(4)]
    HDR_H = 0.52
    row_h = 1.05
    row_y = [GT + HDR_H + GAPY + i * (row_h + GAPY) for i in range(4)]

    # 4) header row — element labels (colored)
    for j, (main, sub, color) in enumerate(COLS):
        b = rrect(slide, col_x[j], GT, col_w, HDR_H, color, None, radius=0.10)
        b.text_frame.word_wrap = True
        tf = b.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), main, F_XB, 12.5, True, WHITE)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        set_run(p2.add_run(), sub, F_R, 8, False, RGBColor(0xE8, 0xEE, 0xF5))

    # corner label above the stage column
    corner = rrect(slide, GL, GT, LBL_W, HDR_H, RGBColor(0x2D, 0x33, 0x38), None, radius=0.10)
    corner.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = corner.text_frame.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    set_run(cp.add_run(), "단계", F_XB, 12.5, True, WHITE)

    # 5) stage rows
    for i, (no, nm, cells) in enumerate(STAGES):
        # stage label cell
        lb = rrect(slide, GL, row_y[i], LBL_W, row_h, STAGE_BG[i], None, radius=0.10)
        lb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = lb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), no, F_XB, 15, True, WHITE)
        p2 = lb.text_frame.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        set_run(p2.add_run(), nm, F_B, 10.5, True, WHITE)
        # element content cells
        for j, txt in enumerate(cells):
            c = rrect(slide, col_x[j], row_y[i], col_w, row_h, CELLBG, CELLLN, radius=0.05)
            c.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            add_seg_para(c.text_frame, parse_bold(txt), 8.6, F_R, INK,
                         align=PP_ALIGN.LEFT, first=True)

    prs.save(SRC)
    print(f"saved {SRC} (slide 18 replaced; shapes now {len(prs.slides[17].shapes)})")


if __name__ == "__main__":
    main()
