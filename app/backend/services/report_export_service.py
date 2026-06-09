"""Binary export builders for final report artifacts."""

from __future__ import annotations

from dataclasses import dataclass
from datetime import date
from io import BytesIO
from pathlib import Path
from typing import Any, Iterable
from xml.sax.saxutils import escape

from docx import Document
from docx.enum.section import WD_SECTION_START
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

from app.backend.schemas.report import (
    FinalReportBundle,
    ReportDocumentModel,
    ReportSection,
    ReportSummaryCard,
)

LABELS = {
    "problem_redefinition": "문제 재정의",
    "target_users": "대상 사용자",
    "prioritized_kpis": "우선 KPI",
    "mvp_scope": "MVP 범위",
    "expected_value": "기대 효과",
    "required_data": "필요 데이터",
    "required_tech_blocks": "필요 기술 블록",
    "constraints": "제약 조건",
    "implementation_order": "구현 순서",
    "verification_plan": "검증 계획",
    "customer_pain": "고객 불편",
    "business_impact": "사업 영향",
    "current_workaround": "현재 우회 방식",
    "success_criteria": "성공 기준",
    "in_scope": "이번 범위",
    "out_of_scope": "제외 범위",
    "acceptance_criteria": "인수 기준",
    "quality_rule": "품질 기준",
    "owner_hint": "담당 힌트",
    "decision_needed": "결정 필요사항",
}

RISK_LABELS = {
    "data": "데이터",
    "technical": "기술",
    "operational": "운영",
    "regulatory": "규제",
    "scope": "범위",
    "other": "기타",
}

SEVERITY_LABELS = {
    "low": "낮음",
    "medium": "중간",
    "high": "높음",
}

PDF_FONT_CANDIDATES = [
    (
        "MilemateNanum",
        "/Library/Fonts/NanumGothic.ttf",
        "/Library/Fonts/NanumGothicBold.ttf",
    ),
    (
        "MilematePretendard",
        "/Users/jaeheemacbook/Library/Fonts/Pretendard-Regular.ttf",
        "/Users/jaeheemacbook/Library/Fonts/Pretendard-Bold.ttf",
    ),
]


@dataclass(frozen=True)
class ReportExport:
    content: bytes
    filename: str
    media_type: str


@dataclass(frozen=True)
class ReportExportMeta:
    """Document metadata that lets a planner route the brief through internal review."""

    document_id: str = ""
    version: str = "1.0"
    created_on: str = ""
    prepared_by: str = "기획팀"


class UnsupportedReportExportError(ValueError):
    """Raised when a requested report export format is not supported."""


def _format_label(value: str) -> str:
    return LABELS.get(value, value.replace("_", " ").title())


def _as_string_list(value: Any) -> list[str]:
    if isinstance(value, list):
        return [str(item) for item in value if str(item)]
    if value:
        return [str(value)]
    return []


def _display_value(value: Any) -> str:
    if isinstance(value, list):
        return "\n".join(_display_value(item) for item in value)
    if isinstance(value, dict):
        return ", ".join(
            f"{_format_label(str(key))}: {_display_value(item)}"
            for key, item in value.items()
        )
    if value is None:
        return "없음"
    return str(value)


def _rows_from_dict(data: dict[str, Any]) -> list[dict[str, str]]:
    return [
        {"구분": _format_label(str(key)), "내용": _display_value(value)}
        for key, value in data.items()
        if value not in (None, "", [])
    ]


def _rows_from_records(records: Iterable[Any]) -> list[dict[str, str]]:
    rows: list[dict[str, str]] = []
    for record in records:
        if hasattr(record, "model_dump"):
            record = record.model_dump(mode="json")
        if isinstance(record, dict):
            rows.append(
                {
                    _format_label(str(key)): _display_value(value)
                    for key, value in record.items()
                    if value not in (None, "", [])
                }
            )
        elif record:
            rows.append({"항목": str(record)})
    return rows


def _pdf_font_names() -> tuple[str, str]:
    for font_name, regular_path, bold_path in PDF_FONT_CANDIDATES:
        if Path(regular_path).exists() and Path(bold_path).exists():
            regular = f"{font_name}Regular"
            bold = f"{font_name}Bold"
            if regular not in pdfmetrics.getRegisteredFontNames():
                pdfmetrics.registerFont(TTFont(regular, regular_path))
            if bold not in pdfmetrics.getRegisteredFontNames():
                pdfmetrics.registerFont(TTFont(bold, bold_path))
            return regular, bold
    return "Helvetica", "Helvetica-Bold"


class ReportExportService:
    """Create business-facing final report exports from a report bundle."""

    def build_document_model(
        self,
        report: FinalReportBundle | dict[str, Any],
        meta: ReportExportMeta | None = None,
    ) -> ReportDocumentModel:
        bundle = (
            report
            if isinstance(report, FinalReportBundle)
            else FinalReportBundle.model_validate(report)
        )
        meta = meta or ReportExportMeta()
        quality = bundle.prd_quality
        review_status = "검토 완료" if quality.status == "ready" else "추가 검토 필요"
        packet = bundle.prd_report
        summary = packet.one_page_summary or bundle.planner_report.problem_redefinition
        risk_tone = "risk" if any(item.severity == "high" for item in bundle.risks) else "warning"
        scope_count = (
            len(packet.scope.in_scope)
            if packet.scope
            else len(bundle.planner_report.mvp_scope)
        )
        cards = [
            ReportSummaryCard(
                title="종합 의견",
                value="기획서 초안 준비",
                detail=summary[:120],
                tone="decision",
            ),
            ReportSummaryCard(
                title="추진 판단",
                value="MVP 검토 가능",
                detail=f"이번 범위 {scope_count}개 항목을 기준으로 개발팀 검토가 가능합니다.",
                tone="ready",
            ),
            ReportSummaryCard(
                title="핵심 리스크",
                value=f"{len(bundle.risks)}건",
                detail="데이터/운영/규제 리스크는 회의 전 확인이 필요합니다.",
                tone=risk_tone,
            ),
            ReportSummaryCard(
                title="다음 회의 안건",
                value=f"{len(packet.decision_agenda)}건",
                detail="범위, 데이터 소유자, 검증 기준을 확정합니다.",
                tone="technical",
            ),
        ]
        sections = [
            ReportSection(
                title="문제 정의",
                body=packet.problem.customer_pain or bundle.planner_report.problem_redefinition,
                items=packet.problem.success_criteria or bundle.planner_report.prioritized_kpis,
                rows=_rows_from_dict(packet.problem.model_dump(mode="json")),
                tone="decision",
            ),
            ReportSection(
                title="대상 사용자와 기대 효과",
                body="비개발 기획자가 아이디어를 설명 가능한 기획서 구조로 바꾸는 데 집중합니다.",
                items=bundle.planner_report.expected_value,
                rows=_rows_from_records(packet.personas),
                tone="ready",
            ),
            ReportSection(
                title="MVP 범위",
                body="이번 단계에서는 실제 배포에 앞서 의사결정 가능한 범위와 보류 조건을 확정합니다.",
                items=packet.scope.in_scope or bundle.planner_report.mvp_scope,
                rows=[
                    {"구분": "이번 범위", "내용": "\n".join(packet.scope.in_scope)},
                    {"구분": "제외 범위", "내용": "\n".join(packet.scope.out_of_scope)},
                ],
                tone="decision",
            ),
            ReportSection(
                title="화면과 운영 정책",
                body="운영자가 실제로 확인할 화면, 정책, 빈 상태와 오류 상태를 개발 전 검토합니다.",
                rows=_rows_from_records(packet.screens[:3] + packet.policies[:3]),
                tone="neutral",
            ),
            ReportSection(
                title="KPI와 데이터 요구사항",
                body="개발 착수 전에 기준값, 목표값, 원천 데이터, 최신성 기준을 확인합니다.",
                rows=_rows_from_records(packet.metrics[:4] + packet.data_requirements[:4]),
                tone="technical",
            ),
            ReportSection(
                title="개발팀 확인사항",
                body="기획서가 코드 작업으로 넘어가기 전에 확인해야 할 전달사항입니다.",
                items=packet.developer_handoff or bundle.engineer_report.verification_plan,
                rows=_rows_from_records(packet.implementation_slices),
                tone="technical",
            ),
        ]
        if bundle.risks:
            sections.append(
                ReportSection(
                    title="리스크 및 대응방안",
                    body=(
                        "리스크는 추진 반대 근거가 아니라, 다음 회의에서 "
                        "확정해야 할 검증 항목입니다."
                    ),
                    rows=[
                        {
                            "구분": RISK_LABELS.get(item.category, item.category),
                            "수준": SEVERITY_LABELS.get(item.severity, item.severity),
                            "내용": item.description,
                            "대응": item.mitigation or "",
                        }
                        for item in bundle.risks
                    ],
                    tone="risk",
                )
            )
        if bundle.decision_log:
            sections.append(
                ReportSection(
                    title="의사결정 이력",
                    body="승인/보류/되돌림 판단 근거가 남아 검토 회의에서 설명할 수 있습니다.",
                    rows=[
                        {
                            "결정 항목": item.item,
                            "상태": item.status,
                            "근거": item.rationale or "",
                        }
                        for item in bundle.decision_log
                    ],
                    tone="decision",
                )
            )
        if bundle.citations:
            sections.append(
                ReportSection(
                    title="참고자료",
                    rows=[
                        {
                            "자료": item.title,
                            "종류": item.source_type,
                            "위치": item.locator,
                            "활용": item.relevance_note,
                        }
                        for item in bundle.citations
                    ],
                    tone="neutral",
                )
            )
        return ReportDocumentModel(
            title="Milemate 최종 기획서",
            subtitle="비개발 기획자의 아이디어를 기획서와 개발팀 확인사항으로 정리한 기술 기획서",
            audience="사업/운영 의사결정권자, 개발 리드",
            summary=summary,
            document_id=meta.document_id,
            version=meta.version,
            created_on=meta.created_on or date.today().isoformat(),
            prepared_by=meta.prepared_by,
            review_status=review_status,
            review_score=quality.score,
            review_findings=quality.findings,
            cards=cards,
            sections=sections,
            risks=bundle.risks,
            decisions=bundle.decision_log,
            citations=bundle.citations,
        )

    def export(
        self,
        report: FinalReportBundle | dict[str, Any],
        export_format: str,
        meta: ReportExportMeta | None = None,
    ) -> ReportExport:
        document = self.build_document_model(report, meta)
        if export_format == "docx":
            return ReportExport(
                content=self.to_docx(document),
                filename="milemate-final-planning-brief.docx",
                media_type=(
                    "application/vnd.openxmlformats-officedocument."
                    "wordprocessingml.document"
                ),
            )
        if export_format == "pdf":
            return ReportExport(
                content=self.to_pdf(document),
                filename="milemate-final-planning-brief.pdf",
                media_type="application/pdf",
            )
        if export_format == "pptx":
            return ReportExport(
                content=self.to_pptx(document),
                filename="milemate-final-presentation-deck.pptx",
                media_type=(
                    "application/vnd.openxmlformats-officedocument."
                    "presentationml.presentation"
                ),
            )
        raise UnsupportedReportExportError(f"unsupported report export: {export_format}")

    def to_docx(self, document_model: ReportDocumentModel) -> bytes:
        doc = Document()
        section = doc.sections[0]
        section.start_type = WD_SECTION_START.NEW_PAGE
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)

        normal = doc.styles["Normal"]
        normal.font.name = "Calibri"
        normal.font.size = Pt(11)
        normal.paragraph_format.space_after = Pt(6)
        normal.paragraph_format.line_spacing = 1.1

        title_style = doc.styles["Title"]
        title_style.font.name = "Calibri"
        title_style.font.size = Pt(24)
        title_style.font.color.rgb = RGBColor(11, 37, 69)

        heading_1 = doc.styles["Heading 1"]
        heading_1.font.name = "Calibri"
        heading_1.font.size = Pt(16)
        heading_1.font.color.rgb = RGBColor(46, 116, 181)
        heading_1.paragraph_format.space_before = Pt(16)
        heading_1.paragraph_format.space_after = Pt(8)

        heading_2 = doc.styles["Heading 2"]
        heading_2.font.name = "Calibri"
        heading_2.font.size = Pt(13)
        heading_2.font.color.rgb = RGBColor(46, 116, 181)
        heading_2.paragraph_format.space_before = Pt(12)
        heading_2.paragraph_format.space_after = Pt(6)

        doc.add_paragraph(document_model.title, style="Title")
        subtitle = doc.add_paragraph(document_model.subtitle)
        subtitle.runs[0].font.color.rgb = RGBColor(87, 96, 106)
        self._add_docx_meta_table(doc, document_model)
        doc.add_paragraph(f"요약: {document_model.summary}")

        doc.add_heading("핵심 판단", level=1)
        card_table = doc.add_table(rows=1, cols=4)
        card_table.style = "Light Shading Accent 1"
        for idx, card in enumerate(document_model.cards):
            cell = card_table.rows[0].cells[idx]
            cell.text = ""
            title = cell.paragraphs[0]
            run = title.add_run(card.title)
            run.bold = True
            run.font.size = Pt(9)
            value = cell.add_paragraph(card.value)
            value.runs[0].bold = True
            value.runs[0].font.size = Pt(12)
            cell.add_paragraph(card.detail)

        for section_item in document_model.sections:
            doc.add_heading(section_item.title, level=1)
            if section_item.body:
                doc.add_paragraph(section_item.body)
            for item in section_item.items:
                doc.add_paragraph(item, style="List Bullet")
            if section_item.rows:
                self._add_docx_table(doc, section_item.rows)

        footer = section.footer.paragraphs[0]
        footer.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        footer.add_run("Milemate 기술 기획서")

        buffer = BytesIO()
        doc.save(buffer)
        return buffer.getvalue()

    def _add_docx_meta_table(self, doc: Document, model: ReportDocumentModel) -> None:
        meta_rows = [
            ("문서번호", model.document_id or "—", "작성일", model.created_on or "—"),
            ("버전", model.version or "—", "작성", model.prepared_by or "—"),
        ]
        table = doc.add_table(rows=0, cols=4)
        table.style = "Table Grid"
        for left_label, left_value, right_label, right_value in meta_rows:
            cells = table.add_row().cells
            cells[0].text = left_label
            cells[1].text = left_value
            cells[2].text = right_label
            cells[3].text = right_value
            for idx in (0, 2):
                for paragraph in cells[idx].paragraphs:
                    for run in paragraph.runs:
                        run.bold = True
        doc.add_paragraph(f"보고 대상: {model.audience}")
        review = doc.add_paragraph(self._review_line(model))
        if model.review_findings:
            review.add_run(f"  ·  점검 필요: {', '.join(model.review_findings[:3])}")
        doc.add_paragraph("")

    def _add_docx_table(self, doc: Document, rows: list[dict[str, str]]) -> None:
        if not rows:
            return
        headers = list(dict.fromkeys(key for row in rows for key in row.keys()))
        table = doc.add_table(rows=1, cols=len(headers))
        table.style = "Table Grid"
        header_cells = table.rows[0].cells
        for idx, header in enumerate(headers):
            header_cells[idx].text = header
            for paragraph in header_cells[idx].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        for row in rows:
            cells = table.add_row().cells
            for idx, header in enumerate(headers):
                cells[idx].text = str(row.get(header, ""))
        doc.add_paragraph("")

    def to_pdf(self, document_model: ReportDocumentModel) -> bytes:
        regular_font, bold_font = _pdf_font_names()
        buffer = BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=A4,
            leftMargin=18 * mm,
            rightMargin=18 * mm,
            topMargin=16 * mm,
            bottomMargin=16 * mm,
            title=document_model.title,
        )
        styles = {
            "title": ParagraphStyle(
                "MilemateTitle",
                fontName=bold_font,
                fontSize=18,
                leading=24,
                textColor=colors.HexColor("#0b2545"),
                spaceAfter=8,
            ),
            "subtitle": ParagraphStyle(
                "MilemateSubtitle",
                fontName=regular_font,
                fontSize=10,
                leading=15,
                textColor=colors.HexColor("#57606a"),
                spaceAfter=8,
            ),
            "heading": ParagraphStyle(
                "MilemateHeading",
                fontName=bold_font,
                fontSize=12.5,
                leading=18,
                textColor=colors.HexColor("#0969da"),
                spaceBefore=9,
                spaceAfter=5,
            ),
            "body": ParagraphStyle(
                "MilemateBody",
                fontName=regular_font,
                fontSize=9.5,
                leading=14,
                spaceAfter=4,
            ),
            "bullet": ParagraphStyle(
                "MilemateBullet",
                fontName=regular_font,
                fontSize=9.2,
                leading=13,
                leftIndent=8 * mm,
                firstLineIndent=-3 * mm,
                spaceAfter=3,
            ),
        }
        story: list[Any] = [
            Paragraph(escape(document_model.title), styles["title"]),
            Paragraph(escape(document_model.subtitle), styles["subtitle"]),
            Paragraph(escape(self._meta_line(document_model)), styles["subtitle"]),
            Paragraph(escape(f"보고 대상: {document_model.audience}"), styles["body"]),
            Paragraph(escape(self._review_line(document_model)), styles["body"]),
            Spacer(1, 3 * mm),
            self._pdf_cards_table(document_model.cards, styles),
            Spacer(1, 4 * mm),
        ]
        for section_item in document_model.sections:
            story.append(Paragraph(escape(section_item.title), styles["heading"]))
            if section_item.body:
                story.append(Paragraph(escape(section_item.body), styles["body"]))
            for item in section_item.items:
                story.append(Paragraph(f"• {escape(item)}", styles["bullet"]))
            if section_item.rows:
                story.append(self._pdf_rows_table(section_item.rows, styles))
                story.append(Spacer(1, 3 * mm))

        def draw_footer(canvas, doc_obj) -> None:
            canvas.saveState()
            canvas.setFont(regular_font, 8)
            canvas.drawString(18 * mm, 10 * mm, "Milemate 기술 기획서")
            canvas.drawRightString(192 * mm, 10 * mm, f"{doc_obj.page}쪽")
            canvas.restoreState()

        doc.build(story, onFirstPage=draw_footer, onLaterPages=draw_footer)
        return buffer.getvalue()

    def _pdf_cards_table(
        self,
        cards: list[ReportSummaryCard],
        styles: dict[str, ParagraphStyle],
    ) -> Table:
        data = [
            [
                Paragraph(
                    f"<b>{escape(card.title)}</b><br/>{escape(card.value)}<br/>"
                    f"<font color='#57606a'>{escape(card.detail)}</font>",
                    styles["body"],
                )
                for card in cards
            ]
        ]
        table = Table(data, colWidths=[43.5 * mm] * 4)
        table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#eef6ff")),
                    ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor("#b6d4fe")),
                    ("INNERGRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#d8dee4")),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 7),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 7),
                    ("TOPPADDING", (0, 0), (-1, -1), 7),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
                ]
            )
        )
        return table

    def _pdf_rows_table(
        self,
        rows: list[dict[str, str]],
        styles: dict[str, ParagraphStyle],
    ) -> Table:
        headers = list(dict.fromkeys(key for row in rows for key in row.keys()))
        data = [[Paragraph(f"<b>{escape(header)}</b>", styles["body"]) for header in headers]]
        for row in rows[:12]:
            data.append(
                [
                    Paragraph(escape(str(row.get(header, ""))), styles["body"])
                    for header in headers
                ]
            )
        width = 174 * mm
        col_width = width / max(len(headers), 1)
        table = Table(data, colWidths=[col_width] * len(headers), repeatRows=1)
        table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#f2f4f7")),
                    ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#d0d7de")),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 5),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 5),
                    ("TOPPADDING", (0, 0), (-1, -1), 5),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
                ]
            )
        )
        return table

    def to_pptx(self, document_model: ReportDocumentModel) -> bytes:
        prs = Presentation()
        prs.slide_width = PptInches(13.333)
        prs.slide_height = PptInches(7.5)
        self._add_title_slide(prs, document_model)
        self._add_summary_slide(prs, document_model)
        self._add_section_slide(prs, "문제와 MVP 범위", document_model.sections[:3])
        self._add_risk_slide(prs, document_model)
        self._add_section_slide(prs, "개발팀 확인사항", document_model.sections[4:6])
        buffer = BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def _add_title_slide(self, prs: Presentation, document_model: ReportDocumentModel) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_band(slide, "#0b2545", 0, 0, 13.333, 1.1)
        self._add_textbox(slide, "Milemate", 0.55, 0.32, 2.0, 0.3, 14, "#ffffff", bold=True)
        self._add_textbox(
            slide,
            document_model.title,
            0.7,
            2.15,
            8.5,
            0.8,
            34,
            "#0b2545",
            bold=True,
        )
        self._add_textbox(slide, document_model.subtitle, 0.72, 3.05, 9.5, 0.6, 18, "#425466")
        self._add_textbox(
            slide,
            self._meta_line(document_model),
            0.72,
            3.72,
            11.8,
            0.35,
            12,
            "#57606a",
        )
        self._add_textbox(
            slide,
            f"보고 대상: {document_model.audience}",
            0.72,
            4.18,
            8.5,
            0.4,
            13,
            "#57606a",
        )
        self._add_textbox(
            slide,
            self._review_line(document_model),
            0.72,
            4.62,
            9.5,
            0.4,
            13,
            "#1a7f37",
            bold=True,
        )
        self._add_band(slide, "#eef6ff", 0.72, 5.15, 11.9, 1.05)
        self._add_textbox(slide, document_model.summary[:210], 1.0, 5.38, 11.2, 0.6, 15, "#24292f")

    def _add_summary_slide(self, prs: Presentation, document_model: ReportDocumentModel) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "핵심 판단")
        for idx, card in enumerate(document_model.cards):
            left = 0.7 + idx * 3.15
            self._add_band(slide, self._tone_color(card.tone), left, 1.45, 2.8, 3.95)
            self._add_textbox(
                slide,
                card.title,
                left + 0.18,
                1.68,
                2.4,
                0.35,
                13,
                "#ffffff",
                bold=True,
            )
            self._add_textbox(
                slide,
                card.value,
                left + 0.18,
                2.15,
                2.4,
                0.45,
                20,
                "#ffffff",
                bold=True,
            )
            self._add_textbox(slide, card.detail[:130], left + 0.18, 2.85, 2.35, 1.4, 12, "#ffffff")

    def _add_section_slide(
        self,
        prs: Presentation,
        title: str,
        sections: list[ReportSection],
    ) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, title)
        for idx, section_item in enumerate(sections[:3]):
            left = 0.72 + idx * 4.08
            self._add_band(slide, "#f6f8fa", left, 1.32, 3.72, 5.3)
            self._add_textbox(
                slide,
                section_item.title,
                left + 0.22,
                1.55,
                3.2,
                0.35,
                15,
                "#0969da",
                bold=True,
            )
            body = section_item.body[:170]
            self._add_textbox(slide, body, left + 0.22, 2.08, 3.25, 0.95, 11, "#24292f")
            items = section_item.items[:4]
            if not items and section_item.rows:
                items = [
                    " / ".join(str(value) for value in row.values() if value)[:90]
                    for row in section_item.rows[:4]
                ]
            self._add_bullets(slide, items, left + 0.28, 3.2, 3.2, 2.1)

    def _add_risk_slide(self, prs: Presentation, document_model: ReportDocumentModel) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "리스크와 의사결정")
        risks = document_model.risks[:4]
        decisions = document_model.decisions[:4]
        self._add_band(slide, "#fff8c5", 0.72, 1.32, 6.05, 5.3)
        self._add_textbox(slide, "핵심 리스크", 1.0, 1.55, 4.0, 0.35, 16, "#7a5a00", bold=True)
        self._add_bullets(
            slide,
            [
                f"{RISK_LABELS.get(item.category, item.category)} · "
                f"{SEVERITY_LABELS.get(item.severity, item.severity)}: {item.description}"
                for item in risks
            ]
            or ["기록된 리스크가 없습니다."],
            1.05,
            2.08,
            5.4,
            3.9,
        )
        self._add_band(slide, "#eefbea", 7.05, 1.32, 5.55, 5.3)
        self._add_textbox(slide, "다음 결정", 7.33, 1.55, 4.0, 0.35, 16, "#1a7f37", bold=True)
        self._add_bullets(
            slide,
            [f"{item.status}: {item.item}" for item in decisions]
            or ["최종 회의에서 범위와 데이터 소유자를 확정합니다."],
            7.38,
            2.08,
            5.0,
            3.9,
        )

    def _add_slide_title(self, slide: Any, title: str) -> None:
        self._add_band(slide, "#0b2545", 0, 0, 13.333, 0.72)
        self._add_textbox(slide, title, 0.72, 0.18, 9.5, 0.32, 18, "#ffffff", bold=True)

    def _add_band(
        self,
        slide: Any,
        hex_color: str,
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            PptInches(left),
            PptInches(top),
            PptInches(width),
            PptInches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._ppt_rgb(hex_color)
        shape.line.color.rgb = self._ppt_rgb(hex_color)

    def _add_textbox(
        self,
        slide: Any,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        size: int,
        color: str,
        bold: bool = False,
    ) -> None:
        textbox = slide.shapes.add_textbox(
            PptInches(left),
            PptInches(top),
            PptInches(width),
            PptInches(height),
        )
        frame = textbox.text_frame
        frame.word_wrap = True
        frame.margin_left = PptInches(0.02)
        frame.margin_right = PptInches(0.02)
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = text
        run.font.name = "Apple SD Gothic Neo"
        run.font.size = PptPt(size)
        run.font.bold = bold
        run.font.color.rgb = self._ppt_rgb(color)

    def _add_bullets(
        self,
        slide: Any,
        items: list[str],
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        textbox = slide.shapes.add_textbox(
            PptInches(left),
            PptInches(top),
            PptInches(width),
            PptInches(height),
        )
        frame = textbox.text_frame
        frame.word_wrap = True
        frame.clear()
        for idx, item in enumerate(items[:5]):
            paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            paragraph.text = f"- {item[:150]}"
            paragraph.level = 0
            paragraph.font.name = "Apple SD Gothic Neo"
            paragraph.font.size = PptPt(11)
            paragraph.font.color.rgb = self._ppt_rgb("#24292f")

    @staticmethod
    def _meta_line(model: ReportDocumentModel) -> str:
        parts: list[str] = []
        if model.document_id:
            parts.append(f"문서번호 {model.document_id}")
        if model.created_on:
            parts.append(f"작성일 {model.created_on}")
        if model.version:
            parts.append(f"버전 {model.version}")
        if model.prepared_by:
            parts.append(f"작성 {model.prepared_by}")
        return "   ·   ".join(parts)

    @staticmethod
    def _review_line(model: ReportDocumentModel) -> str:
        line = f"검토 상태: {model.review_status or '미상'}"
        if model.review_score:
            line += f" (품질 점수 {model.review_score}/100)"
        return line

    @staticmethod
    def _tone_color(tone: str) -> str:
        return {
            "decision": "#0969da",
            "ready": "#1a7f37",
            "warning": "#bf8700",
            "risk": "#cf222e",
            "technical": "#57606a",
        }.get(tone, "#57606a")

    @staticmethod
    def _ppt_rgb(hex_color: str):
        from pptx.dml.color import RGBColor as PptRGBColor

        clean = hex_color.lstrip("#")
        return PptRGBColor(
            int(clean[0:2], 16),
            int(clean[2:4], 16),
            int(clean[4:6], 16),
        )
